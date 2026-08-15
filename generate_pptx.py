import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

with open("src/assets/i18n/en.json", "r", encoding="utf-8") as f:
    t = json.load(f)

NAVY = RGBColor(0x06, 0x0E, 0x1E)
LIGHT_NAVY = RGBColor(0x0D, 0x1F, 0x3C)
LIGHT_SLATE = RGBColor(0xA8, 0xB2, 0xD1)
LIGHTEST_SLATE = RGBColor(0xCC, 0xD6, 0xF6)
WHITE = RGBColor(0xE6, 0xF1, 0xFF)
GREEN = RGBColor(0x1B, 0xCE, 0xDF)
SLATE = RGBColor(0x88, 0x92, 0xB0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=LIGHTEST_SLATE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_accent_bar(slide, left, top, width=Inches(1.5), height=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    return shape


def add_section_number(slide, number, top=Inches(0.4)):
    add_textbox(slide, Inches(0.6), top, Inches(1), Inches(0.5),
                number, font_size=14, color=GREEN, bold=True, font_name="Consolas")


def add_title(slide, title, top=Inches(0.35)):
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.8),
                title, font_size=32, color=WHITE, bold=True)


def add_lead(slide, text, top=Inches(1.3)):
    add_textbox(slide, Inches(0.6), top, Inches(12), Inches(1),
                text, font_size=17, color=LIGHT_SLATE)


def add_card(slide, left, top, width, height, title, items):
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_NAVY
    card.line.color.rgb = RGBColor(0x1A, 0x33, 0x56)
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # Title
    add_textbox(slide, left + Inches(0.25), top + Inches(0.15),
                width - Inches(0.5), Inches(0.4),
                title, font_size=16, color=WHITE, bold=True)

    # Items
    txBox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.6),
                                     width - Inches(0.5), height - Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▹  {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_SLATE
        p.space_after = Pt(6)
        p.line_spacing = 1.2


# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0.6), Inches(2.2), Inches(2))
add_textbox(slide, Inches(0.6), Inches(2.5), Inches(12), Inches(1.5),
            t["School.Hero.Title"], font_size=44, color=WHITE, bold=True)
add_textbox(slide, Inches(0.6), Inches(4.0), Inches(11), Inches(1.2),
            t["School.Hero.Subtitle"], font_size=18, color=LIGHT_SLATE)
add_textbox(slide, Inches(0.6), Inches(5.5), Inches(4), Inches(0.5),
            t["School.Hero.Tag"], font_size=14, color=GREEN, bold=True, font_name="Consolas")

# ── Slide 2: Intro ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, "01.")
add_title(slide, t["School.Intro.Title"])
add_lead(slide, t["School.Intro.Description"])
add_accent_bar(slide, Inches(0.6), Inches(2.6), Inches(1.5))
add_textbox(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(0.5),
            t["School.Intro.HighlightTitle"], font_size=20, color=WHITE, bold=True)
add_textbox(slide, Inches(0.6), Inches(3.4), Inches(12), Inches(0.8),
            f"• {t['School.Intro.Highlight1']}", font_size=16, color=LIGHT_SLATE)
add_textbox(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.8),
            f"• {t['School.Intro.Highlight2']}", font_size=16, color=LIGHT_SLATE)
add_textbox(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.6),
            t["School.Intro.HighlightNote"], font_size=15, color=GREEN, bold=True)

# ── Slide 3: Social Network ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, "02.")
add_title(slide, t["School.Social.Title"])
add_lead(slide, t["School.Social.Lead"])

cards = [
    (t["School.Social.Feed.Title"], [t["School.Social.Feed.Item1"], t["School.Social.Feed.Item2"], t["School.Social.Feed.Item3"]]),
    (t["School.Social.Stories.Title"], [t["School.Social.Stories.Item1"], t["School.Social.Stories.Item2"], t["School.Social.Stories.Item3"]]),
    (t["School.Social.Engagement.Title"], [t["School.Social.Engagement.Item1"], t["School.Social.Engagement.Item2"]]),
    (t["School.Social.Targeting.Title"], [t["School.Social.Targeting.Item1"], t["School.Social.Targeting.Item2"], t["School.Social.Targeting.Item3"]]),
    (t["School.Social.Safety.Title"], [t["School.Social.Safety.Item1"], t["School.Social.Safety.Item2"], t["School.Social.Safety.Item3"]]),
]
card_w = Inches(2.35)
card_h = Inches(3.2)
gap = Inches(0.15)
start_left = Inches(0.5)
for i, (title, items) in enumerate(cards):
    left = start_left + (card_w + gap) * i
    add_card(slide, left, Inches(2.5), card_w, card_h, title, items)

add_textbox(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(0.8),
            t["School.Social.Result"], font_size=15, color=GREEN)

# ── Slide 4: Sawaed AI ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, "03.")
add_title(slide, t["School.AI.Title"])
add_lead(slide, t["School.AI.Lead"])

# Left: Questions
add_textbox(slide, Inches(0.6), Inches(2.5), Inches(6), Inches(0.5),
            t["School.AI.Questions.Title"], font_size=18, color=WHITE, bold=True)
questions = [t["School.AI.Questions.Q1"], t["School.AI.Questions.Q2"], t["School.AI.Questions.Q3"],
             t["School.AI.Questions.Q4"], t["School.AI.Questions.Q5"]]
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(6), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
for i, q in enumerate(questions):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"▹  {q}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_SLATE
    p.space_after = Pt(8)
note_p = tf.add_paragraph()
note_p.text = t["School.AI.Questions.Note"]
note_p.font.size = Pt(13)
note_p.font.color.rgb = GREEN
note_p.font.italic = True
note_p.space_before = Pt(10)

# Right: Data sources
add_textbox(slide, Inches(7.0), Inches(2.5), Inches(6), Inches(0.5),
            t["School.AI.Data.Title"], font_size=18, color=WHITE, bold=True)
data_items = [t[f"School.AI.Data.Item{i}"] for i in range(1, 8)]
txBox2 = slide.shapes.add_textbox(Inches(7.0), Inches(3.0), Inches(6), Inches(3))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, item in enumerate(data_items):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = f"▹  {item}"
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_SLATE
    p.space_after = Pt(6)

add_textbox(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.8),
            t["School.AI.Result"], font_size=15, color=GREEN)

# ── Slide 5: Feature Catalog (part 1) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, "04.")
add_title(slide, t["School.CatalogTitle"])
add_lead(slide, t["School.CatalogLead"])

catalog = t["School.Catalog"]
half = (len(catalog) + 1) // 2
card_w = Inches(4.0)
card_h = Inches(3.8)
gap = Inches(0.2)
for i, cat in enumerate(catalog[:half]):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + (card_w + gap) * col
    top = Inches(2.3) + (card_h + gap) * row
    items = cat["items"][:5]
    add_card(slide, left, top, card_w, card_h, f"{cat['icon']}  {cat['title']}", items)

# ── Slide 6: Feature Catalog (part 2) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, "04.")
add_title(slide, t["School.CatalogTitle"] + " (cont.)")

for i, cat in enumerate(catalog[half:]):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + (card_w + gap) * col
    top = Inches(0.8) + (card_h + gap) * row
    items = cat["items"][:5]
    add_card(slide, left, top, card_w, card_h, f"{cat['icon']}  {cat['title']}", items)

# ── Slide 7: Problems Solved ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, "05.")
add_title(slide, t["School.ProblemsTitle"])

problems = t["School.Problems"]
row_h = Inches(0.55)
for i, row in enumerate(problems[:7]):
    top = Inches(1.8) + row_h * i
    # Problem
    add_textbox(slide, Inches(0.6), top, Inches(5.8), row_h,
                f"✗  {row['problem']}", font_size=12, color=SLATE)
    # Solution
    add_textbox(slide, Inches(6.6), top, Inches(6.3), row_h,
                f"✓  {row['solution']}", font_size=12, color=GREEN)

# ── Slide 8: Benefits ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, "06.")
add_title(slide, t["School.BenefitsTitle"])

benefits = t["School.Benefits"]
card_w = Inches(4.0)
card_h = Inches(2.3)
gap = Inches(0.2)
for i, b in enumerate(benefits):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + (card_w + gap) * col
    top = Inches(1.8) + (card_h + gap) * row
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_NAVY
    card.line.color.rgb = RGBColor(0x1A, 0x33, 0x56)
    card.line.width = Pt(1)
    card.shadow.inherit = False
    add_textbox(slide, left + Inches(0.25), top + Inches(0.15),
                card_w - Inches(0.5), Inches(0.5),
                f"{b['icon']}  {b['title']}", font_size=16, color=WHITE, bold=True)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.7),
                card_w - Inches(0.5), card_h - Inches(0.9),
                b["text"], font_size=13, color=LIGHT_SLATE)

# ── Slide 9: Comparison ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_number(slide, "07.")
add_title(slide, t["School.ComparisonTitle"])

comparison = t["School.Comparison"]
# Headers
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(3), Inches(0.4),
            "Aspect", font_size=14, color=GREEN, bold=True, font_name="Consolas")
add_textbox(slide, Inches(4.0), Inches(1.6), Inches(4), Inches(0.4),
            t["School.ComparisonHeader2"], font_size=14, color=SLATE, bold=True)
add_textbox(slide, Inches(8.5), Inches(1.6), Inches(4.5), Inches(0.4),
            t["School.ComparisonHeader3"], font_size=14, color=GREEN, bold=True)

for i, row in enumerate(comparison):
    top = Inches(2.2) + Inches(0.7) * i
    add_textbox(slide, Inches(0.6), top, Inches(3), Inches(0.6),
                row["aspect"], font_size=14, color=WHITE, bold=True)
    add_textbox(slide, Inches(4.0), top, Inches(4), Inches(0.6),
                row["typical"], font_size=13, color=SLATE)
    add_textbox(slide, Inches(8.5), top, Inches(4.5), Inches(0.6),
                row["sawaed"], font_size=13, color=GREEN)

# ── Slide 10: CTA ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0.6), Inches(2.5), Inches(2))
add_textbox(slide, Inches(0.6), Inches(2.8), Inches(12), Inches(1),
            t["School.CTA.Title"], font_size=40, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(1),
            t["School.CTA.Text"], font_size=18, color=LIGHT_SLATE,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.5),
            "info@sawaedelm.com", font_size=20, color=GREEN, bold=True,
            alignment=PP_ALIGN.CENTER)

output_path = "Sawaed_Elm_School_Platform.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
